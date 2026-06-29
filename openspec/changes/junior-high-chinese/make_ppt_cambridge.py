# -*- coding: utf-8 -*-
"""產生「我所知道的康橋」教學 PPT（15 張，課文全讀 7 段各佔一頁，20pt）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

# ── 色盤 ─────────────────────────────────────────────
C_DEEP   = RGBColor(0x1B, 0x3A, 0x5C)
C_MEADOW = RGBColor(0x2E, 0x6A, 0x4F)
C_SKY    = RGBColor(0x3A, 0x7A, 0xAA)
C_GOLD   = RGBColor(0xB5, 0x8C, 0x28)
C_CREAM  = RGBColor(0xFA, 0xF7, 0xEE)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x4A, 0x50, 0x60)
C_PANEL  = RGBColor(0xE8, 0xF3, 0xF8)
C_LIGHT  = RGBColor(0xD8, 0xEF, 0xE4)
C_COVER  = RGBColor(0x0C, 0x1E, 0x30)
C_STEEL  = RGBColor(0x7A, 0xA0, 0xC0)
C_LGRAY  = RGBColor(0xDD, 0xE8, 0xF0)

CT = Inches(1.22)
CH = H - CT - Inches(0.1)   # ≈ 6.18"

# ── 工具函式 ─────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txbox(slide, text, l, t, w, h,
          size=Pt(16), bold=False, color=C_DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def multi_para(slide, lines, l, t, w, h,
               size=Pt(16), color=C_DARK, align=PP_ALIGN.LEFT,
               bold=False, italic=False, spacing=Pt(8)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = spacing
        r = p.add_run()
        r.text = line
        r.font.size  = size
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.italic = italic
    return box

def strip_top(slide, color=C_DEEP, h=Inches(0.16)):
    rect(slide, 0, 0, W, h, fill=color)

def strip_bot(slide, color=C_MEADOW, h=Inches(0.1)):
    rect(slide, 0, H - h, W, h, fill=color)

def slide_title(slide, title):
    rect(slide, Inches(0.52), Inches(0.35), Inches(0.08), Inches(0.58), fill=C_GOLD)
    txbox(slide, title,
          Inches(0.72), Inches(0.37),
          Inches(11.5), Inches(0.62),
          size=Pt(24), bold=True, color=C_DEEP)

def badge(slide, text, l, t, w=Inches(1.4), h=Inches(0.42),
          fill=C_SKY, tc=C_WHITE, size=Pt(14)):
    rect(slide, l, t, w, h, fill=fill)
    txbox(slide, text, l, t, w, h,
          size=size, bold=True, color=tc, align=PP_ALIGN.CENTER)

def full_read(slide, label, text, accent=None):
    """全讀頁面：白底大框 + 左側色條 + 段落標籤 + 20pt 內文"""
    ac = accent or C_SKY
    # 白底內容框
    rect(slide, Inches(0.50), CT, Inches(12.33), CH,
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    # 左側色條
    rect(slide, Inches(0.50), CT, Inches(0.16), CH, fill=ac)
    # 段落標籤
    badge(slide, label, Inches(0.80), CT + Inches(0.14),
          w=Inches(1.25), h=Inches(0.44), fill=ac, size=Pt(14))
    # 內文（20pt）
    txbox(slide, text,
          Inches(0.82), CT + Inches(0.70),
          Inches(12.00), CH - Inches(0.80),
          size=Pt(20), color=C_DARK)

# ── 課文文字 ─────────────────────────────────────────

P1 = ("靜極了，這朝來水溶溶的大道，只遠處牛奶車的鈴聲，點綴這個週遭的沉默。"
      "順著這大道走去，走到盡頭，再轉入林子裏的小徑，往煙霧濃密處走去，"
      "頭頂是交枝的榆蔭，透露著漠楞楞的曙色，再往前走去，走盡這林子，"
      "當前是平坦的原野，望見了村舍，初青的麥田，"
      "更遠三兩個饅形的小山掩住了一條通道。"
      "天邊是霧茫茫的，尖尖的黑影是近村的教寺。聽，那曉鐘和緩的清音。")

P2 = ("這一帶是此邦中部的平原，地形像是海面的清波，默沈沈的起伏；"
      "山嶺是望不見的，有的是常青的草原與沃腴的田壤。"
      "登那土阜上望去，康橋只是一帶茂林，擁戴著幾處娉婷的尖閣。"
      "嫵媚的康河也望不見蹤跡，你只能循著那錦帶似的林木想像那一流清淺。")

P3 = ("村舍與樹林是這地盤上的棋子，有村舍有佳蔭，有佳蔭處有村舍。"
      "這早起是看炊煙的時辰；朝霧漸漸的升起，揭開了這灰蒼蒼的天幕，"
      "（最好是微霰後的光景）遠近的炊煙，成絲的，成縷的，成捲的，"
      "輕快的，遲重的，濃灰的，淡青的，慘白的，"
      "在靜定的朝氣裏漸漸的上騰，漸漸的不見，"
      "彷彿是朝來人們的祈躊，參差的翳入了天聽。")

P4 = ("朝陽是難得見的，這初春的天氣，但它來時是起早人莫大的愉快。"
      "頃刻間這田野添深了顏色，一層輕紗似的金粉糝上了這草，"
      "這樹，這通道，這莊舍。"
      "頃刻間這周遭瀰漫了清晨富麗的溫柔。"
      "頃刻間你的心懷也分潤了白天誕生的光榮。"
      "「春！」這勝利的晴空彷彿在你的耳邊私語。"
      "「春！」妳那快活的靈魂也彷彿在那裏回響。")

P5 = ("伺候著河上的風光，這春來一天有一天的消息。"
      "關心石上的苔痕，關心敗草裏的鮮花，關心這水流的緩急，"
      "關心水草的滋長，關心天上的雲霞，關心新來的鳥語。"
      "怯怜怜的小雪球是探春信的小使。鈴蘭與香草是歡喜的初聲。"
      "窈窕的蓮馨，玲瓏的石水仙，愛熱鬧的克羅克斯，"
      "耐辛苦的蒲公英與雛菊──"
      "這時候春光已是爛縵在人間，更不須殷勤問訊。")

P6 = ("瑰麗的春天。這是你野遊的時期，可愛的路政，"
      "這裏不比中國，那一處不是坦蕩蕩的大道？"
      "徒步是一個愉快，但騎自轉車是一個更大的愉快。"
      "在康橋騎車是普遍的技術；婦人，稚子，老翁，一致享受這雙輪舞的快樂。"
      "（在康橋聽說自轉車是不怕人偷的，就為人人都自己有車，沒人要偷。）"
      "任你選一個方向，任你上一條通道，順著這帶草味的和風，"
      "放輪遠去，保管你這半天的逍遙是你靈性的補劑。"
      "這道上有的是清蔭與美草，隨地都可以供你休憩。"
      "你如愛花，這裏多的是錦繡似的草原。"
      "你如愛鳥，這裏多的是巧囀的鳴禽。"
      "你如愛兒童，這鄉間到處是可親的稚子。"
      "你如愛人情，這裏多的是不嫌遠客的鄉人，"
      "你到處可以「掛單」借宿，有酪漿與嫩薯供你飽餐，有奪目的鮮果恣你嘗新。"
      "你如愛酒，這鄉間每「望」都為妳儲有上好的新釀，"
      "黑啤如太濃，蘋果酒、薑酒都是供你解渴潤肺的。"
      "……帶一卷書，走十里路，選一塊清靜地，"
      "看天，聽鳥，讀書，倦了時，和身在草綿綿處尋夢去"
      "──你能想像更適情更適性的消遣嗎？")

P7 = ("陸放翁有一聯詩句：「傳呼快馬迎新月，卻上輕輿趁晚涼；」這是做地方官的風流。"
      "我在康橋時雖沒馬騎，沒轎子坐，卻也有我的風流；"
      "我常常在夕陽西曬時騎了車迎著天邊扁大的日頭直追。"
      "日頭是追不到的，我沒有夸父的荒誕，但晚景的溫存卻被我這樣偷嘗了不少。"
      "只說看夕陽，我們平常只知道登山或是臨海，"
      "但實際只須遼闊的天際，平地上的晚霞有時也是一樣的神奇。"
      "有一次我趕到一個地方，手把著一家村莊的籬巴，"
      "隔著一大田的麥浪，看西天的變幻。"
      "有一次是正衝著一條寬廣的大道，過來一大群羊，放草歸來的，"
      "偌大的太陽在它們後背放射著萬縷的金輝。"
      "天上卻是烏青青的，只賸這不可逼視的威光中的一條大路，一群生物！"
      "我心頭頓時感著神異性的壓迫，我真的跪下了，對著這冉冉漸翳的金光。"
      "再有一次是更不可忘的奇景，那是臨著一大片望不到頭的草原，"
      "滿開著豔紅的罌粟，在青草裏亭亭的像是萬盞的金燈，"
      "陽光從褐色雲裏斜著過來，幻成一種異樣的紫色，透明似的不可逼視，"
      "霎那間在我迷眩了的視覺中，這草田變成了……不說也罷，說來你們也是不信的！")

# ══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ─────────────────────────────────────────────────────
# S1  封面
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)

rect(s, Inches(9.1), 0, Inches(4.23), H, fill=RGBColor(0x08, 0x14, 0x20))
rect(s, Inches(9.05), 0, Inches(0.06), H, fill=C_SKY)
rect(s, Inches(0.55), Inches(2.65), Inches(8.1), Inches(0.05), fill=C_GOLD)

txbox(s, "我所知道的康橋",
      Inches(0.55), Inches(1.05), Inches(8.1), Inches(1.5),
      size=Pt(50), bold=True, color=C_WHITE)
txbox(s, "徐志摩",
      Inches(0.55), Inches(2.75), Inches(4.5), Inches(0.6),
      size=Pt(22), color=C_SKY, bold=True)
txbox(s, "選自《我所知道的康橋》，第三冊第9課",
      Inches(0.55), Inches(3.38), Inches(8.1), Inches(0.5),
      size=Pt(14), color=C_STEEL, italic=True)
badge(s, "八年級上學期",
      Inches(0.55), Inches(4.12), w=Inches(1.9), fill=C_GOLD, tc=C_DARK)
multi_para(s, ["英倫", "田野", "印象"],
           Inches(9.3), Inches(2.2), Inches(3.7), Inches(2.8),
           size=Pt(30), bold=True,
           color=RGBColor(0x2A, 0x50, 0x40),
           align=PP_ALIGN.CENTER, spacing=Pt(4))
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_GOLD)

# ─────────────────────────────────────────────────────
# S2  作者介紹
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "作者介紹")

rect(s, Inches(0.55), CT, Inches(3.9), Inches(1.05), fill=C_DEEP)
txbox(s, "徐志摩",
      Inches(0.55), CT + Inches(0.08), Inches(3.9), Inches(0.55),
      size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Xu Zhimo  1897–1931",
      Inches(0.55), CT + Inches(0.60), Inches(3.9), Inches(0.38),
      size=Pt(12), color=C_STEEL, align=PP_ALIGN.CENTER, italic=True)

for i, tag in enumerate(["浙江海寧", "劍橋留學", "新月詩派"]):
    badge(s, tag, Inches(0.55 + i * 1.38), CT + Inches(1.22),
          w=Inches(1.25), fill=C_SKY)

items = [
    "▸  1921–1922 年赴英國劍橋大學（King's College）留學，"
    "一年的康橋生涯深深影響了他對自然美的感受與文學創作風格",
    "▸  與胡適、聞一多共同創立新月詩社，提倡詩的「音樂美、建築美、繪畫美」，"
    "是中國現代新詩運動的核心人物",
    "▸  散文語言抒情優美、感官豐富，代表作《我所知道的康橋》《再別康橋》"
    "至今被視為中國現代散文與詩歌的雙璧，1931 年飛機失事英年早逝",
]
multi_para(s, items,
           Inches(0.65), CT + Inches(1.82), Inches(12.0), Inches(4.0),
           size=Pt(16), color=C_GRAY, spacing=Pt(14))

# ─────────────────────────────────────────────────────
# S3  背景知識
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "文化背景：劍橋與徐志摩")

for i, (icon, head, body) in enumerate([
    ("🏛", "留學康橋（1921–1922）",
     "劍橋大學（Cambridge University）是英國最古老的大學之一，"
     "「康橋」是 Cambridge 的中文音譯。徐志摩在此旁聽文學課程，"
     "每日漫步劍河兩岸，深受英國田園風光感化。"),
    ("🌿", "散文的感官美學",
     "課文以視覺（炊煙顏色、晨光金粉、豔紅罌粟）、"
     "聽覺（鈴聲、晨鐘、鳥語）、嗅覺（草味和風）三感交織，"
     "語言抒情濃郁，被稱為「詩化散文」的典範。"),
    ("✒", "名句的誕生",
     "「帶一卷書，走十里路，選一塊清靜地，看天，聽鳥，讀書，"
     "倦了時，和身在草綿綿處尋夢去」，"
     "這句話精煉了作者在康橋的理想生活哲學，至今廣為流傳。"),
]):
    col = i * Inches(4.3) + Inches(0.50)
    rect(s, col, CT + Inches(0.05), Inches(4.05), Inches(5.5),
         fill=C_WHITE, line_color=C_SKY, line_w=Pt(1.5))
    txbox(s, icon, col + Inches(0.15), CT + Inches(0.12),
          Inches(0.8), Inches(0.7), size=Pt(28), align=PP_ALIGN.CENTER)
    txbox(s, head, col + Inches(0.15), CT + Inches(0.85),
          Inches(3.75), Inches(0.55),
          size=Pt(17), bold=True, color=C_DEEP)
    txbox(s, body, col + Inches(0.15), CT + Inches(1.48),
          Inches(3.75), Inches(3.9),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S4–S10  課文全讀（各段各頁，20pt）
# ─────────────────────────────────────────────────────
slides_data = [
    ("課文全文（一）", "清晨漫步", P1, C_SKY),
    ("課文全文（二）", "康橋平原", P2, C_SKY),
    ("課文全文（三）", "炊煙百態", P3, C_DEEP),
    ("課文全文（四）", "晨光春聲", P4, C_MEADOW),
    ("課文全文（五）", "春日河景", P5, C_MEADOW),
    ("課文全文（六）", "騎車郊遊", P6, C_SKY),
    ("課文全文（七）", "追夕陽",   P7, C_GOLD),
]

for title, label, text, accent in slides_data:
    s = new_slide(prs)
    bg(s, C_CREAM)
    strip_top(s); strip_bot(s)
    slide_title(s, title)
    full_read(s, label, text, accent=accent)

# ─────────────────────────────────────────────────────
# S11  重要詞彙（六詞 2×3 緊湊版）
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "重要詞彙")

vocab_all = [
    # (字, 注音, 釋義, 出處句, 色)
    ("溶溶", "ㄖㄨㄥˊ ㄖㄨㄥˊ",
     "水流廣大浩淼、緩緩漫延的樣子，形容晨曦下大道水氣瀰漫的意境",
     "這朝來水溶溶的大道", C_SKY),
    ("娉婷", "ㄆㄧㄥ ㄊㄧㄥˊ",
     "姿態修長優美，文中以擬人化描寫教堂尖塔的秀麗挺拔",
     "擁戴著幾處娉婷的尖閣", C_SKY),
    ("嫵媚", "ㄨˇ ㄇㄟˋ",
     "姿容柔美動人，這裡以擬人化賦予康河女性般的柔情風韻",
     "嫵媚的康河也望不見蹤跡", C_SKY),
    ("翳入", "ㄧˋ ㄖㄨˋ",
     "遮蔽而漸漸消散融入，指炊煙緩緩升騰最終融入天際",
     "參差的翳入了天聽", C_MEADOW),
    ("逍遙", "ㄒㄧㄠ ㄧㄠˊ",
     "自由自在、悠遊無拘的狀態，形容騎車漫遊時心靈的舒暢",
     "你這半天的逍遙是你靈性的補劑", C_MEADOW),
    ("罌粟", "ㄧㄥ ㄙㄨˋ",
     "一年生草本植物，開豔紅大花，文中形容草原如萬盞金燈的壯麗景致",
     "滿開著豔紅的罌粟，亭亭的像是萬盞的金燈", C_MEADOW),
]

COLS = 3
ROW_H = (CH - Inches(0.08)) / 2   # 每列高度，約 3.05"
COL_W = (W - Inches(0.90)) / COLS  # 每欄寬度，約 4.14"
GAP   = Inches(0.06)

for idx, (word, phonetic, defn, ex, hc) in enumerate(vocab_all):
    row = idx // COLS
    col = idx % COLS
    l = Inches(0.45) + col * (COL_W + GAP)
    t = CT + Inches(0.04) + row * (ROW_H + GAP)
    w = COL_W
    h = ROW_H

    rect(s, l, t, w, h, fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, w, Inches(0.80), fill=hc)

    # 字
    txbox(s, word,
          l + Inches(0.12), t + Inches(0.06),
          Inches(1.8), Inches(0.65),
          size=Pt(26), bold=True, color=C_WHITE)
    # 注音
    txbox(s, phonetic,
          l + Inches(1.95), t + Inches(0.22),
          Inches(2.0), Inches(0.40),
          size=Pt(12), color=RGBColor(0xD0, 0xE8, 0xFF), italic=True)
    # 分隔線
    rect(s, l + Inches(0.12), t + Inches(0.88),
         w - Inches(0.24), Inches(0.03), fill=C_LGRAY)
    # 釋義
    txbox(s, defn,
          l + Inches(0.12), t + Inches(0.96),
          w - Inches(0.24), Inches(1.30),
          size=Pt(14), color=C_DARK)
    # 例句底框
    rect(s, l + Inches(0.12), t + h - Inches(0.68),
         w - Inches(0.24), Inches(0.58), fill=C_PANEL)
    txbox(s, "例：" + ex,
          l + Inches(0.18), t + h - Inches(0.65),
          w - Inches(0.36), Inches(0.54),
          size=Pt(12), color=C_MEADOW, italic=True)

# ─────────────────────────────────────────────────────
# S12  主題分析
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "主題分析")

for i, (head, hcolor, items) in enumerate([
    ("感官的饗宴", C_SKY,
     ["視覺：炊煙色彩、金粉晨光、豔紅罌粟",
      "聽覺：牛奶車鈴聲、曉鐘清音、鳥語",
      "嗅覺：帶草味的和風",
      "三感交織，呈現自然的多層次美感",
      "這是「詩化散文」的核心技法"]),
    ("自然的情懷", C_MEADOW,
     ["以擬人化賦予自然生命",
      "「娉婷的尖閣」「嫵媚的康河」",
      "炊煙如「人們的祈躊」",
      "自然不是被觀賞的客體",
      "而是與人平等共情的存在"]),
    ("生活的美學", C_GOLD,
     ["「帶一卷書，走十里路，",
      " 看天，聽鳥，讀書，",
      " 和身在草綿綿處尋夢去」",
      "傳遞「讀書人應親近自然」",
      "是全文最廣為人知的名句"]),
]):
    col = i * Inches(4.25) + Inches(0.45)
    rect(s, col, CT + Inches(0.05), Inches(4.0), Inches(5.5),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, col, CT + Inches(0.05), Inches(4.0), Inches(0.72), fill=hcolor)
    txbox(s, head, col + Inches(0.12), CT + Inches(0.10),
          Inches(3.75), Inches(0.62),
          size=Pt(19), bold=True, color=C_WHITE)
    multi_para(s, ["▸  " + it for it in items],
               col + Inches(0.18), CT + Inches(0.92),
               Inches(3.7), Inches(4.5),
               size=Pt(15), color=C_GRAY, spacing=Pt(10))

# ─────────────────────────────────────────────────────
# S13  核心金句
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)
rect(s, 0, 0, W, H, fill=RGBColor(0x08, 0x14, 0x20))

txbox(s, "核心金句",
      Inches(0.6), Inches(0.28), Inches(5.0), Inches(0.55),
      size=Pt(16), bold=True, color=C_SKY)
rect(s, Inches(0.6), Inches(0.80), Inches(8.0), Inches(0.05), fill=C_GOLD)

for i, (q, qc) in enumerate([
    ("「頃刻間你的心懷也分潤了白天誕生的光榮。\n"
     "  「春！」妳那快活的靈魂也彷彿在那裏回響。」",
     C_WHITE),
    ("「帶一卷書，走十里路，選一塊清靜地，\n"
     "  看天，聽鳥，讀書，倦了時，\n"
     "  和身在草綿綿處尋夢去──」",
     RGBColor(0xA0, 0xD8, 0xC0)),
    ("「我真的跪下了，對著這冉冉漸翳的金光。」",
     C_GOLD),
]):
    t = Inches(1.05) + i * Inches(2.02)
    rect(s, Inches(0.55), t + Inches(0.12),
         Inches(0.06), Inches(1.65),
         fill=C_GOLD if i == 2 else C_SKY)
    txbox(s, q, Inches(0.75), t,
          Inches(12.3), Inches(1.85),
          size=Pt(19 if i == 1 else 18),
          bold=(i == 2), color=qc, italic=True)

# ─────────────────────────────────────────────────────
# S14  課堂討論
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課堂討論")

for i, (qno, qtxt) in enumerate([
    ("Q1", "作者以哪些感官（視、聽、嗅）描寫康橋的晨景？\n試各舉課文中一個例子。"),
    ("Q2", "描寫炊煙時用了哪些形容詞？\n這樣密集的形容詞有什麼效果？"),
    ("Q3", "「帶一卷書，走十里路……和身在草綿綿處尋夢去」——\n你認為作者最嚮往的生活是什麼樣子？"),
    ("Q4", "描寫夕陽的三幅畫面中，你最喜歡哪一幅？\n說說感受，並說明原因。"),
]):
    row, col_i = i // 2, i % 2
    l = Inches(0.45) + col_i * Inches(6.4)
    t = CT + Inches(0.08) + row * Inches(2.9)
    rect(s, l, t, Inches(6.1), Inches(2.72),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, Inches(1.0), Inches(2.72), fill=C_DEEP)
    txbox(s, qno, l + Inches(0.05), t + Inches(0.95),
          Inches(0.88), Inches(0.82),
          size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txbox(s, qtxt, l + Inches(1.12), t + Inches(0.20),
          Inches(4.82), Inches(2.32),
          size=Pt(15.5), color=C_DARK)

# ─────────────────────────────────────────────────────
# S15  總結
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)
rect(s, 0, 0, Inches(5.5), H, fill=RGBColor(0x08, 0x14, 0x20))
rect(s, Inches(5.5), 0, Inches(0.06), H, fill=C_SKY)

txbox(s, "學習重點總結",
      Inches(0.5), Inches(0.6), Inches(4.5), Inches(0.6),
      size=Pt(16), bold=True, color=C_SKY)
rect(s, Inches(0.5), Inches(1.12), Inches(4.5), Inches(0.05), fill=C_GOLD)

multi_para(s,
    ["徐志摩", "浙江海寧人", "劍橋大學（1921–1922）", "",
     "新月詩派", "詩化散文代表", "感官美學"],
    Inches(0.55), Inches(1.28), Inches(4.5), Inches(5.5),
    size=Pt(16), color=RGBColor(0xB0, 0xD0, 0xC8), spacing=Pt(8))

for i, (head, body) in enumerate([
    ("感官交織的自然描寫",   "視覺、聽覺、嗅覺三感並用，建構詩化境界"),
    ("擬人化的情感投射",     "娉婷的尖閣、嫵媚的康河，自然即情感"),
    ("生活美學的宣言",       "「帶一卷書，走十里路」是全文最廣傳的名句"),
    ("三幅夕陽畫，震撼收尾", "「我真的跪下了」將自然之美推至情感頂點"),
]):
    t = CT - Inches(0.04) + i * Inches(1.52)
    rect(s, Inches(5.85), t, Inches(7.0), Inches(1.40),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, Inches(5.85), t, Inches(0.28), Inches(1.40), fill=C_MEADOW)
    txbox(s, head, Inches(6.22), t + Inches(0.12),
          Inches(6.45), Inches(0.52),
          size=Pt(17), bold=True, color=C_DEEP)
    txbox(s, body, Inches(6.22), t + Inches(0.65),
          Inches(6.45), Inches(0.62),
          size=Pt(14.5), color=C_GRAY)

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_GOLD)

# ── 儲存 ─────────────────────────────────────────────
OUT = r"C:\Users\yuchi\OneDrive\Desktop\教學日誌\我所知道的康橋.pptx"
prs.save(OUT)
print(f"已儲存：{OUT}")
