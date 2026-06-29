# -*- coding: utf-8 -*-
"""產生「水神的指引」教學 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 尺寸與色盤 ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

C_DEEP   = RGBColor(0x1B, 0x3A, 0x4B)
C_FOREST = RGBColor(0x2D, 0x6A, 0x4F)
C_ACCENT = RGBColor(0x52, 0xB7, 0x88)
C_GOLD   = RGBColor(0xE9, 0xC4, 0x6A)
C_CREAM  = RGBColor(0xF8, 0xF5, 0xEE)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY   = RGBColor(0x4A, 0x4A, 0x5A)
C_LGRAY  = RGBColor(0xE8, 0xE8, 0xF0)
C_LIGHT_GREEN = RGBColor(0xD8, 0xF3, 0xDC)
C_DARK2  = RGBColor(0x14, 0x2D, 0x3C)

# ── 基礎函式 ─────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.background() if fill is None else (sh.fill.solid(), sh.fill.fore_color.__setattr__('rgb', fill))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txbox(slide, text, l, t, w, h,
          size=Pt(16), bold=False, color=C_DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = size
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def multi_para(slide, lines, l, t, w, h,
               size=Pt(16), color=C_DARK, align=PP_ALIGN.LEFT,
               bold=False, italic=False, spacing=Pt(8),
               bold_first=False):
    """多段落文字框"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = spacing
        r = p.add_run()
        r.text = line
        r.font.size  = size
        r.font.color.rgb = color
        r.font.bold  = bold_first if (line == lines[0]) else bold
        r.font.italic = italic
    return box

def strip_top(slide, color=C_DEEP, h=Inches(0.16)):
    rect(slide, 0, 0, W, h, fill=color)

def strip_bot(slide, color=C_FOREST, h=Inches(0.1)):
    rect(slide, 0, H - h, W, h, fill=color)

def side_accent(slide, t=Inches(0.38), h=Inches(0.55), color=C_GOLD):
    rect(slide, Inches(0.52), t, Inches(0.08), h, fill=color)

def slide_title(slide, title, t=Inches(0.35)):
    side_accent(slide, t=t, h=Inches(0.58))
    txbox(slide, title,
          Inches(0.72), t + Inches(0.02),
          Inches(11.5), Inches(0.62),
          size=Pt(24), bold=True, color=C_DEEP)

def badge(slide, text, l, t, w=Inches(1.4), h=Inches(0.38),
          fill=C_ACCENT, tc=C_WHITE, size=Pt(13)):
    rect(slide, l, t, w, h, fill=fill)
    txbox(slide, text, l, t, w, h,
          size=size, bold=True, color=tc, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ─────────────────────────────────────────────────────
# S1  封面
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_DEEP)

rect(s, Inches(9.0), 0, Inches(4.33), H, fill=C_DARK2)
rect(s, Inches(8.95), 0, Inches(0.05), H, fill=C_ACCENT)
rect(s, Inches(0.55), Inches(2.72), Inches(8.0), Inches(0.055), fill=C_GOLD)

txbox(s, "水神的指引",
      Inches(0.55), Inches(1.15), Inches(8.0), Inches(1.5),
      size=Pt(56), bold=True, color=C_WHITE)

txbox(s, "亞榮隆‧撒可努",
      Inches(0.55), Inches(2.82), Inches(5.0), Inches(0.6),
      size=Pt(20), color=C_ACCENT, bold=True)

txbox(s, "選自《山豬‧飛鼠‧撒可努》，INK 印刻出版",
      Inches(0.55), Inches(3.42), Inches(8.0), Inches(0.5),
      size=Pt(14), color=RGBColor(0xA0, 0xC4, 0xB8), italic=True)

badge(s, "九年級下學期",
      Inches(0.55), Inches(4.15), w=Inches(1.9), fill=C_GOLD, tc=C_DARK)

multi_para(s,
    ["人與自然", "共生共息"],
    Inches(9.2), Inches(2.3), Inches(3.7), Inches(2.8),
    size=Pt(32), bold=True, color=RGBColor(0x2D, 0x5A, 0x50),
    align=PP_ALIGN.CENTER, spacing=Pt(4))

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_GOLD)

# ─────────────────────────────────────────────────────
# S2  作者介紹
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "作者介紹")

rect(s, Inches(0.55), Inches(1.15), Inches(3.9), Inches(1.1), fill=C_FOREST)
txbox(s, "亞榮隆‧撒可努",
      Inches(0.55), Inches(1.18), Inches(3.9), Inches(0.58),
      size=Pt(21), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Ahronglong Sakinu",
      Inches(0.55), Inches(1.72), Inches(3.9), Inches(0.42),
      size=Pt(13), color=RGBColor(0xC0, 0xE8, 0xD8),
      align=PP_ALIGN.CENTER, italic=True)

for i, tag in enumerate(["排灣族", "屏東獅子鄉", "原住民文學"]):
    badge(s, tag, Inches(0.55 + i * 1.38), Inches(2.42), w=Inches(1.25), fill=C_ACCENT)

items = [
    "▸  以自身在部落的成長經歷為素材，用中文書寫原住民狩獵文化與生命哲學",
    "▸  代表作《山豬‧飛鼠‧撒可努》細膩描繪父子之間的知識傳承，榮獲多項文學獎",
    "▸  首位以中文書寫並獲主流文壇肯定的排灣族作家，被譽為原住民文學代表人物",
]
multi_para(s, items,
           Inches(0.65), Inches(3.05), Inches(12.0), Inches(3.8),
           size=Pt(17), color=C_GRAY, spacing=Pt(10))

# ─────────────────────────────────────────────────────
# S3  背景知識
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "文化背景：排灣族與山林教育")

for i, (icon, head, body) in enumerate([
    ("🏔", "狩獵即教育",
     "在排灣族傳統中，父親帶兒子入山狩獵不只是技術傳授，\n更是完整的生命教育——如何閱讀自然訊息、在天地間找到定位。"),
    ("🌿", "取捨之道",
     "取走需要的，剩餘的還給山神與水神。\n這套永續生態觀早於現代環保概念數百年，根植於部落文化之中。"),
    ("💧", "水神信仰",
     "文中的「水神」不是迷信，而是對大自然力量的尊重與詮釋。\n動物是水神的牛羊，蘊含人與自然共享山林的倫理觀念。"),
]):
    col = i * Inches(4.3) + Inches(0.55)
    rect(s, col, Inches(1.2), Inches(4.05), Inches(5.4),
         fill=C_WHITE, line_color=C_ACCENT, line_w=Pt(1.5))
    txbox(s, icon, col + Inches(0.15), Inches(1.25), Inches(0.8), Inches(0.7),
          size=Pt(30), align=PP_ALIGN.CENTER)
    txbox(s, head, col + Inches(0.15), Inches(1.98), Inches(3.75), Inches(0.52),
          size=Pt(18), bold=True, color=C_DEEP)
    txbox(s, body, col + Inches(0.15), Inches(2.52), Inches(3.75), Inches(3.8),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S4  課文（一）段 1–2
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課文全文（一）")

rect(s, Inches(0.55), Inches(1.18), Inches(12.2), Inches(5.55),
     fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))

badge(s, "第一段", Inches(0.65), Inches(1.3), w=Inches(1.1), fill=C_FOREST)
txbox(s, "父親說：「山裡的動物，是水神放牧在山上的牛羊，你要學會和牠們分享，牠們才會在你的森林裡長大、繁衍。」",
      Inches(1.88), Inches(1.3), Inches(10.7), Inches(1.15),
      size=Pt(16), color=C_DARK)

rect(s, Inches(0.65), Inches(2.6), Inches(12.0), Inches(0.04), fill=C_LGRAY)

badge(s, "第二段", Inches(0.65), Inches(2.72), w=Inches(1.1), fill=C_FOREST)
txbox(s, "有一次，我們在森林裡遇見一場大雨，雨水像神明的淚水一樣，從天空傾瀉而下。父親帶領我走在沒有路的路徑上，尋找一處可以遮風避雨的岩洞。那時我還小，感到非常恐懼，但父親卻說：「不要怕，那是水神在引路。」",
      Inches(1.88), Inches(2.72), Inches(10.7), Inches(1.7),
      size=Pt(16), color=C_DARK)

# ─────────────────────────────────────────────────────
# S5  課文（二）段 3–4
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課文全文（二）")

rect(s, Inches(0.55), Inches(1.18), Inches(12.2), Inches(5.55),
     fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))

badge(s, "第三段", Inches(0.65), Inches(1.3), w=Inches(1.1), fill=C_FOREST)
txbox(s, "我們在岩洞裡生火取暖，雨水滴落在洞口，發出清脆的節奏。父親指著遠方說：「你看，那些動物也正在躲雨，牠們在看著我們。你要記得，人和自然是共生的，我們取走我們需要的，剩下的要還給山神和水神。如果你心存感激，水神就會引領你找到出口。」",
      Inches(1.88), Inches(1.3), Inches(10.7), Inches(2.1),
      size=Pt(16), color=C_DARK)

rect(s, Inches(0.65), Inches(3.6), Inches(12.0), Inches(0.04), fill=C_LGRAY)

badge(s, "第四段", Inches(0.65), Inches(3.72), w=Inches(1.1), fill=C_FOREST)
txbox(s, "那天晚上，父親教我如何聽雨的聲音，如何分辨風傳來的訊息。他說，水流動的方向就是生命的脈動，只要你安靜下來，你就能聽見山林的呼吸。雨停後，月光灑在溼漉漉的樹葉上，像是點點星光落在人間。",
      Inches(1.88), Inches(3.72), Inches(10.7), Inches(1.7),
      size=Pt(16), color=C_DARK)

# ─────────────────────────────────────────────────────
# S6  課文（三）段 5
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課文全文（三）")

rect(s, Inches(0.55), Inches(1.18), Inches(12.2), Inches(4.5),
     fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))

badge(s, "第五段", Inches(0.65), Inches(1.3), w=Inches(1.1), fill=C_FOREST)
txbox(s, "父親帶著我走出森林，我發現原本陌生恐懼的山林，在水神的指引下，變得如此親近而神聖。我學會了不再用征服的眼光看待自然，而是學會彎下腰，用謙卑的心去領受大地的恩賜。",
      Inches(1.88), Inches(1.3), Inches(10.7), Inches(2.0),
      size=Pt(17), color=C_DARK)

txbox(s, "——　亞榮隆‧撒可努〈水神的指引〉，選自《山豬‧飛鼠‧撒可努》",
      Inches(2.5), Inches(5.85), Inches(10.0), Inches(0.45),
      size=Pt(13), color=RGBColor(0x88, 0x88, 0x99), italic=True,
      align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────
# S7  段落解析（一）段 1–2
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "段落解析（一）")

for row, (para_label, quote, analysis) in enumerate([
    ("第一段",
     "「山裡的動物，是水神放牧在山上的牛羊……」",
     "以「放牧」比喻野生動物與水神的關係，奠定全文的神話視角。\n動物不是可隨意獵殺的獵物，而是受神明庇護的生命，人與牠們是平等的共居者。"),
    ("第二段",
     "「不要怕，那是水神在引路。」",
     "面對未知的恐懼，父親用神話語言給孩子安全感。\n水神不是抽象的神明，而是大自然本身——雨水有方向，山林有意志，讀懂就能找到出路。"),
]):
    t = Inches(1.3) + row * Inches(2.85)
    rect(s, Inches(0.55), t, Inches(12.2), Inches(2.65),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    badge(s, para_label, Inches(0.65), t + Inches(0.12),
          w=Inches(1.1), fill=C_DEEP)
    txbox(s, quote, Inches(1.88), t + Inches(0.1), Inches(10.6), Inches(0.65),
          size=Pt(15), color=C_FOREST, italic=True, bold=True)
    rect(s, Inches(1.88), t + Inches(0.82), Inches(10.6), Inches(0.04), fill=C_LGRAY)
    txbox(s, analysis, Inches(1.88), t + Inches(0.94), Inches(10.6), Inches(1.55),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S8  段落解析（二）段 3–4
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "段落解析（二）")

for row, (para_label, quote, analysis) in enumerate([
    ("第三段",
     "「人和自然是共生的，我們取走我們需要的，剩下的要還給山神和水神。」",
     "這是全文最核心的生態倫理宣言。「取」與「還」構成一個閉合的循環，\n體現排灣族永續共生的傳統智慧，與現代環保思想高度吻合。"),
    ("第四段",
     "「水流動的方向就是生命的脈動，只要你安靜下來，你就能聽見山林的呼吸。」",
     "父親將感官教育昇華為生命哲學。「安靜」不只是停下腳步，\n而是放下人類中心的主導欲望，以開放謙虛的態度聆聽大自然。"),
]):
    t = Inches(1.3) + row * Inches(2.85)
    rect(s, Inches(0.55), t, Inches(12.2), Inches(2.65),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    badge(s, para_label, Inches(0.65), t + Inches(0.12),
          w=Inches(1.1), fill=C_DEEP)
    txbox(s, quote, Inches(1.88), t + Inches(0.1), Inches(10.6), Inches(0.65),
          size=Pt(15), color=C_FOREST, italic=True, bold=True)
    rect(s, Inches(1.88), t + Inches(0.82), Inches(10.6), Inches(0.04), fill=C_LGRAY)
    txbox(s, analysis, Inches(1.88), t + Inches(0.94), Inches(10.6), Inches(1.55),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S9  段落解析（三）段 5 + 成長弧線
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "段落解析（三）")

rect(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(2.6),
     fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
badge(s, "第五段", Inches(0.65), Inches(1.42), w=Inches(1.1), fill=C_DEEP)
txbox(s, "「學會彎下腰，用謙卑的心去領受大地的恩賜。」",
      Inches(1.88), Inches(1.4), Inches(10.6), Inches(0.65),
      size=Pt(15), color=C_FOREST, italic=True, bold=True)
rect(s, Inches(1.88), Inches(2.12), Inches(10.6), Inches(0.04), fill=C_LGRAY)
txbox(s, "「彎下腰」與「謙卑」是全文的情感落點——孩子從恐懼到謙遜，從征服視角到領受心態，完成了生命教育的完整弧線。",
      Inches(1.88), Inches(2.24), Inches(10.6), Inches(1.5),
      size=Pt(14.5), color=C_GRAY)

txbox(s, "孩子的成長弧線",
      Inches(0.65), Inches(4.1), Inches(3.5), Inches(0.45),
      size=Pt(16), bold=True, color=C_DEEP)

for i, (icon, label) in enumerate([
    ("😨", "恐懼"), ("👁", "觀察"), ("💡", "領悟"), ("🙏", "謙卑")
]):
    cx = Inches(0.55) + i * Inches(3.0)
    rect(s, cx, Inches(4.65), Inches(2.75), Inches(1.8),
         fill=C_LIGHT_GREEN, line_color=C_ACCENT, line_w=Pt(1.2))
    txbox(s, icon, cx + Inches(0.1), Inches(4.7), Inches(0.8), Inches(0.75),
          size=Pt(26), align=PP_ALIGN.CENTER)
    txbox(s, label, cx + Inches(0.85), Inches(4.88), Inches(1.75), Inches(0.55),
          size=Pt(17), bold=True, color=C_DEEP)
    if i < 3:
        txbox(s, "→", cx + Inches(2.78), Inches(5.0), Inches(0.3), Inches(0.45),
              size=Pt(20), color=C_ACCENT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────
# S10  主題分析
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "主題分析")

themes = [
    ("🌿", "敬畏自然",
     "文中的「水神」不是迷信，而是排灣族對大自然力量的詮釋。\n把動物比作「水神的牛羊」，以神話語言表達生態倫理：\n動物與人共享山林，不是單方面的獵取對象。"),
    ("♻", "人與自然共生",
     "「我們取走我們需要的，剩下的要還給山神和水神」\n體現原住民傳統的永續生態觀，\n取與還形成一個閉合的生態倫理循環。"),
    ("🙏", "謙卑的生命態度",
     "文章最核心的成長弧線：\n孩子從「恐懼」走向「謙卑」，\n從「征服」視角轉為「領受」心態，\n是生命教育最珍貴的結果。"),
]

for i, (icon, head, body) in enumerate(themes):
    col = i * Inches(4.25) + Inches(0.45)
    rect(s, col, Inches(1.22), Inches(4.0), Inches(5.5),
         fill=C_WHITE, line_color=C_ACCENT, line_w=Pt(1.5))
    rect(s, col, Inches(1.22), Inches(4.0), Inches(0.75), fill=C_FOREST)
    txbox(s, icon + "  " + head,
          col + Inches(0.12), Inches(1.25), Inches(3.75), Inches(0.68),
          size=Pt(18), bold=True, color=C_WHITE)
    txbox(s, body, col + Inches(0.15), Inches(2.08), Inches(3.72), Inches(4.4),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S11  文學特色
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "文學特色")

features = [
    ("靈性的自然書寫",
     C_FOREST,
     [
         "將山林描寫為充滿生命與意志的存在",
         "雨水是「神明的淚水」，水流有「生命的脈動」",
         "月光「像是點點星光落在人間」",
         "靈性視角是原住民文學有別於一般自然散文的核心特質",
     ]),
    ("父子對話的教育結構",
     C_DEEP,
     [
         "全文透過父親的引導語推動敘事前進",
         "父親每說一句話，都是一堂課",
         "孩子沒有提問，只有觀察與感悟",
         "沉默式的學習，正是部落教育的傳統形式",
     ]),
]

for i, (head, hcolor, items) in enumerate(features):
    col = i * Inches(6.2) + Inches(0.45)
    rect(s, col, Inches(1.22), Inches(5.95), Inches(5.55),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, col, Inches(1.22), Inches(5.95), Inches(0.72), fill=hcolor)
    txbox(s, head, col + Inches(0.12), Inches(1.26), Inches(5.7), Inches(0.62),
          size=Pt(18), bold=True, color=C_WHITE)
    multi_para(s,
               ["▸  " + it for it in items],
               col + Inches(0.18), Inches(2.08),
               Inches(5.65), Inches(4.5),
               size=Pt(15.5), color=C_GRAY, spacing=Pt(12))

# ─────────────────────────────────────────────────────
# S12  重要詞彙（一）
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "重要詞彙（一）")

vocab1 = [
    ("傾瀉", "ㄑㄧㄥ ㄒㄧㄝˋ",
     "大量液體急速、大幅度地流下或湧出。",
     "雨水像神明的淚水一樣，從天空傾瀉而下"),
    ("共生", "ㄍㄨㄥˋ ㄕㄥ",
     "兩種以上的生命體彼此相互依存、共同生活。",
     "人和自然是共生的"),
    ("謙卑", "ㄑㄧㄢ ㄅㄟ",
     "虛心低調，不自大傲慢，以開放臣服的態度面對自然。",
     "用謙卑的心去領受大地的恩賜"),
]

for i, (word, phonetic, defn, ex) in enumerate(vocab1):
    col = i * Inches(4.25) + Inches(0.45)
    rect(s, col, Inches(1.22), Inches(4.0), Inches(5.5),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, col, Inches(1.22), Inches(4.0), Inches(1.0), fill=C_DEEP)
    txbox(s, word,
          col + Inches(0.12), Inches(1.27), Inches(2.0), Inches(0.82),
          size=Pt(30), bold=True, color=C_WHITE)
    txbox(s, phonetic,
          col + Inches(2.15), Inches(1.46), Inches(1.75), Inches(0.48),
          size=Pt(13), color=RGBColor(0xC0, 0xE0, 0xD8), italic=True)
    rect(s, col + Inches(0.15), Inches(2.38), Inches(3.72), Inches(0.04),
         fill=C_LGRAY)
    txbox(s, defn,
          col + Inches(0.15), Inches(2.5), Inches(3.72), Inches(1.6),
          size=Pt(14.5), color=C_DARK)
    rect(s, col + Inches(0.15), Inches(4.25), Inches(3.72), Inches(0.8),
         fill=C_LIGHT_GREEN)
    txbox(s, "例：" + ex,
          col + Inches(0.22), Inches(4.28), Inches(3.55), Inches(0.72),
          size=Pt(13), color=C_FOREST, italic=True)

# ─────────────────────────────────────────────────────
# S13  重要詞彙（二）
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "重要詞彙（二）")

vocab2 = [
    ("領受", "",
     "接受、承受（多指珍貴或重要的事物）。以感恩之心承納大地所給予的一切。",
     "用謙卑的心去領受大地的恩賜"),
    ("脈動", "ㄇㄞˋ ㄉㄨㄥˋ",
     "原指脈搏的跳動節奏，引申為生命力的展現與流動。",
     "水流動的方向就是生命的脈動"),
    ("放牧", "",
     "讓牲畜在戶外自由吃草覓食。文中以此比擬野生動物受神明庇護，蘊含生態倫理。",
     "山裡的動物，是水神放牧在山上的牛羊"),
]

for i, (word, phonetic, defn, ex) in enumerate(vocab2):
    col = i * Inches(4.25) + Inches(0.45)
    rect(s, col, Inches(1.22), Inches(4.0), Inches(5.5),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, col, Inches(1.22), Inches(4.0), Inches(1.0), fill=C_DEEP)
    txbox(s, word,
          col + Inches(0.12), Inches(1.27), Inches(2.0), Inches(0.82),
          size=Pt(30), bold=True, color=C_WHITE)
    if phonetic:
        txbox(s, phonetic,
              col + Inches(2.15), Inches(1.46), Inches(1.75), Inches(0.48),
              size=Pt(13), color=RGBColor(0xC0, 0xE0, 0xD8), italic=True)
    rect(s, col + Inches(0.15), Inches(2.38), Inches(3.72), Inches(0.04),
         fill=C_LGRAY)
    txbox(s, defn,
          col + Inches(0.15), Inches(2.5), Inches(3.72), Inches(1.6),
          size=Pt(14.5), color=C_DARK)
    rect(s, col + Inches(0.15), Inches(4.25), Inches(3.72), Inches(0.8),
         fill=C_LIGHT_GREEN)
    txbox(s, "例：" + ex,
          col + Inches(0.22), Inches(4.28), Inches(3.55), Inches(0.72),
          size=Pt(13), color=C_FOREST, italic=True)

# ─────────────────────────────────────────────────────
# S14  核心金句
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_DEEP)
rect(s, 0, 0, W, H, fill=C_DARK2)

txbox(s, "核心金句",
      Inches(0.6), Inches(0.3), Inches(5.0), Inches(0.55),
      size=Pt(16), bold=True, color=C_ACCENT)
rect(s, Inches(0.6), Inches(0.82), Inches(7.5), Inches(0.055), fill=C_GOLD)

quotes = [
    ("「人和自然是共生的，我們取走我們需要的，\n  剩下的要還給山神和水神。」", C_WHITE),
    ("「水流動的方向就是生命的脈動，\n  只要你安靜下來，你就能聽見山林的呼吸。」",
     RGBColor(0xA0, 0xD8, 0xBE)),
    ("「學會彎下腰，用謙卑的心去領受大地的恩賜。」",
     C_GOLD),
]

for i, (q, qc) in enumerate(quotes):
    t = Inches(1.05) + i * Inches(1.98)
    rect(s, Inches(0.55), t + Inches(0.1),
         Inches(0.06), Inches(1.55), fill=C_GOLD if i == 2 else C_ACCENT)
    txbox(s, q, Inches(0.75), t, Inches(12.3), Inches(1.8),
          size=Pt(19 if i == 2 else 17), bold=(i == 2),
          color=qc, italic=True)

# ─────────────────────────────────────────────────────
# S15  課堂討論
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課堂討論")

qs = [
    ("Q1", "父親說「那是水神在引路」，你認為這代表什麼意思？\n這是真正的神明，還是另一種表達方式？"),
    ("Q2", "「取走需要的，剩下的要還給山神和水神」——\n這個觀念和現代環保思想有什麼異同？"),
    ("Q3", "文章結尾孩子的心態有什麼改變？\n是什麼讓他從恐懼走向謙卑？"),
    ("Q4", "如果你是那個孩子，你會從父親身上學到什麼？\n這和你日常生活有什麼連結？"),
]

for i, (qno, qtxt) in enumerate(qs):
    row = i // 2
    col_i = i % 2
    l = Inches(0.45) + col_i * Inches(6.4)
    t = Inches(1.28) + row * Inches(2.82)
    rect(s, l, t, Inches(6.1), Inches(2.62),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, Inches(1.0), Inches(2.62), fill=C_FOREST)
    txbox(s, qno, l + Inches(0.05), t + Inches(0.9), Inches(0.88), Inches(0.8),
          size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txbox(s, qtxt, l + Inches(1.12), t + Inches(0.18), Inches(4.82), Inches(2.25),
          size=Pt(15), color=C_DARK)

# ─────────────────────────────────────────────────────
# S16  總結
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_DEEP)
rect(s, Inches(0), 0, Inches(5.5), H, fill=C_DARK2)
rect(s, Inches(5.5), 0, Inches(0.06), H, fill=C_ACCENT)

txbox(s, "學習重點總結",
      Inches(0.5), Inches(0.6), Inches(4.5), Inches(0.6),
      size=Pt(16), bold=True, color=C_ACCENT)
rect(s, Inches(0.5), Inches(1.12), Inches(4.5), Inches(0.055), fill=C_GOLD)

summary_left = [
    "亞榮隆‧撒可努",
    "排灣族作家",
    "原住民文學代表人物",
    "",
    "靈性自然書寫",
    "父子教育敘事結構",
    "生命成長弧線",
]
multi_para(s, summary_left,
           Inches(0.55), Inches(1.28), Inches(4.5), Inches(5.5),
           size=Pt(16), color=RGBColor(0xC0, 0xD8, 0xD0), spacing=Pt(8))

points = [
    ("敬畏自然",   "以神話語言詮釋生態倫理"),
    ("人與自然共生", "取與還的永續生命觀"),
    ("謙卑領受",   "從恐懼走向謙遜的成長弧線"),
    ("原住民智慧",  "數百年前的傳統與現代環保吻合"),
]

for i, (head, body) in enumerate(points):
    t = Inches(1.18) + i * Inches(1.48)
    rect(s, Inches(5.85), t, Inches(7.0), Inches(1.35),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, Inches(5.85), t, Inches(0.28), Inches(1.35), fill=C_FOREST)
    txbox(s, head, Inches(6.22), t + Inches(0.12), Inches(6.45), Inches(0.5),
          size=Pt(17), bold=True, color=C_DEEP)
    txbox(s, body, Inches(6.22), t + Inches(0.62), Inches(6.45), Inches(0.6),
          size=Pt(14.5), color=C_GRAY)

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_GOLD)

# ── 儲存 ─────────────────────────────────────────────
OUT = r"C:\Users\yuchi\OneDrive\Desktop\教學日誌\水神的指引.pptx"
prs.save(OUT)
print(f"已儲存：{OUT}")
