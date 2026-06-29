# -*- coding: utf-8 -*-
"""產生「與宋元思書」教學 PPT（16 張，課文全讀 5 段各佔一頁，20pt）"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

# ── 色盤（富春江水墨風格）─────────────────────────────
C_RIVER    = RGBColor(0x1B, 0x42, 0x5A)   # 深水藍（主色）
C_MOUNTAIN = RGBColor(0x2E, 0x5A, 0x3A)   # 山嵐綠
C_MIST     = RGBColor(0x4A, 0x7A, 0x9B)   # 霧藍
C_AMBER    = RGBColor(0xA8, 0x72, 0x1E)   # 琥珀
C_CREAM    = RGBColor(0xF5, 0xF0, 0xE8)   # 沙米底色
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY     = RGBColor(0x4A, 0x50, 0x60)
C_PANEL    = RGBColor(0xE4, 0xF0, 0xF0)
C_LGRAY    = RGBColor(0xD8, 0xE4, 0xEC)
C_COVER    = RGBColor(0x0A, 0x1C, 0x28)
C_MUTED    = RGBColor(0x7A, 0xA0, 0xB8)

CT = Inches(1.22)
CH = H - CT - Inches(0.1)   # ≈ 6.18"

# ── 工具函式 ─────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txbox(slide, text, l, t, w, h,
          size=Pt(16), bold=False, color=C_DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size    = size
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    return box

def multi_para(slide, lines, l, t, w, h,
               size=Pt(16), color=C_DARK, align=PP_ALIGN.LEFT,
               bold=False, italic=False, spacing=Pt(8)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment    = align
        p.space_before = spacing
        r = p.add_run()
        r.text           = line
        r.font.size      = size
        r.font.color.rgb = color
        r.font.bold      = bold
        r.font.italic    = italic
    return box

def strip_top(slide, color=C_RIVER, h=Inches(0.16)):
    rect(slide, 0, 0, W, h, fill=color)

def strip_bot(slide, color=C_MOUNTAIN, h=Inches(0.1)):
    rect(slide, 0, H - h, W, h, fill=color)

def slide_title(slide, title):
    rect(slide, Inches(0.52), Inches(0.35), Inches(0.08), Inches(0.58), fill=C_AMBER)
    txbox(slide, title,
          Inches(0.72), Inches(0.37),
          Inches(11.5), Inches(0.62),
          size=Pt(24), bold=True, color=C_RIVER)

def badge(slide, text, l, t, w=Inches(1.4), h=Inches(0.42),
          fill=C_MIST, tc=C_WHITE, size=Pt(14)):
    rect(slide, l, t, w, h, fill=fill)
    txbox(slide, text, l, t, w, h,
          size=size, bold=True, color=tc, align=PP_ALIGN.CENTER)

def full_read(slide, label, text, accent=None):
    """全讀頁面：白底大框 + 左側色條 + 段落標籤 + 20pt 內文"""
    ac = accent or C_MIST
    rect(slide, Inches(0.50), CT, Inches(12.33), CH,
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(slide, Inches(0.50), CT, Inches(0.16), CH, fill=ac)
    badge(slide, label, Inches(0.80), CT + Inches(0.14),
          w=Inches(1.25), h=Inches(0.44), fill=ac, size=Pt(14))
    txbox(slide, text,
          Inches(0.82), CT + Inches(0.70),
          Inches(12.00), CH - Inches(0.80),
          size=Pt(20), color=C_DARK)

def full_read_bilingual(slide, label, orig_text, trans_text, accent=None):
    """全讀頁面：原文（20pt）+ 分隔 + 今譯（14pt）"""
    ac = accent or C_MIST
    ORG_H = Inches(1.80)
    SEP_Y = CT + Inches(0.70) + ORG_H + Inches(0.08)
    TRN_Y = SEP_Y + Inches(0.02) + Inches(0.08) + Inches(0.30) + Inches(0.06)
    TRN_H = CT + CH - TRN_Y - Inches(0.08)

    rect(slide, Inches(0.50), CT, Inches(12.33), CH,
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(slide, Inches(0.50), CT, Inches(0.16), CH, fill=ac)
    badge(slide, label, Inches(0.80), CT + Inches(0.14),
          w=Inches(1.25), h=Inches(0.44), fill=ac, size=Pt(14))
    txbox(slide, orig_text,
          Inches(0.82), CT + Inches(0.70),
          Inches(12.00), ORG_H,
          size=Pt(22), color=C_DARK)
    rect(slide, Inches(0.82), SEP_Y, Inches(11.80), Inches(0.02), fill=C_LGRAY)
    badge(slide, "今譯",
          Inches(0.82), SEP_Y + Inches(0.05),
          w=Inches(0.80), h=Inches(0.30), fill=ac, tc=C_WHITE, size=Pt(12))
    txbox(slide, trans_text,
          Inches(0.82), TRN_Y,
          Inches(12.00), TRN_H,
          size=Pt(16), color=C_GRAY)

# ── 課文文字 ─────────────────────────────────────────

P1 = "風煙俱淨，天山共色，從流飄蕩，任意東西。自富陽至桐廬，一百許里，奇山異水，天下獨絕。"

P2 = "水皆縹碧，千丈見底；游魚細石，直視無礙。急湍甚箭，猛浪若奔。"

P3 = "夾峰高山，皆生寒樹。負勢競上，互相軒邈，爭高直指，千百成峰。"

P4 = ("泉水激石，泠泠作響。好鳥相鳴，嚶嚶成韻。"
      "蟬則千轉不窮，猨則百叫無絕。"
      "鳶飛戾天者，望峰息心；經綸世務者，窺谷忘反。")

P5 = "橫柯上蔽，在晝猶昏；疎條交映，有時見日。"

# ── 今譯（簡編本）────────────────────────────────────

T1 = ("風和煙霧都消散了，天空和山巒呈現同一種顏色。"
      "隨著水流飄蕩，任憑船往東往西漂流。"
      "從富陽到桐廬，大約一百多里，奇特的山、不同尋常的水，"
      "是天下獨一無二的絕景。")

T2 = ("水都是青白色的，深達千丈也能清楚看見水底；"
      "游動的魚兒和細小的石子，直視到底毫無阻礙。"
      "湍急的水流比箭還快，洶湧的波浪如同飛奔的馬匹。")

T3 = ("兩岸的高山，都長著令人感到寒意的樹木。"
      "山勢競相向上攀升，互相比試高遠，筆直地爭著聳立高空，"
      "形成了千百個山峰。")

T4 = ("泉水激打著石頭，發出清脆的泠泠水聲。"
      "美麗的鳥兒相互鳴叫，嚶嚶的叫聲和諧成韻。"
      "蟬兒鳴叫個不停，猿猴啼聲不絕於耳。"
      "那些一心追求功名的人，望見這山峰，名利之心便平息了；"
      "那些忙於政務的人，窺見這深幽的山谷，都忘了回去的路。")

T5 = ("橫斜的樹幹在上方遮蔽著，白天也好像黃昏一般昏暗；"
      "稀疏的樹枝相互交錯輝映，偶爾也能見到陽光透射進來。")

# ══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ─────────────────────────────────────────────────────
# S1  封面
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)

rect(s, Inches(9.1),  0, Inches(4.23), H, fill=RGBColor(0x07, 0x10, 0x1A))
rect(s, Inches(9.05), 0, Inches(0.06), H, fill=C_MIST)
rect(s, Inches(0.55), Inches(2.65), Inches(8.1), Inches(0.05), fill=C_AMBER)

txbox(s, "與宋元思書",
      Inches(0.55), Inches(1.05), Inches(8.1), Inches(1.5),
      size=Pt(50), bold=True, color=C_WHITE)
txbox(s, "吳均",
      Inches(0.55), Inches(2.75), Inches(4.5), Inches(0.6),
      size=Pt(22), color=C_MIST, bold=True)
txbox(s, "南朝梁‧書信體古文，第五冊第5課",
      Inches(0.55), Inches(3.38), Inches(8.1), Inches(0.5),
      size=Pt(14), color=C_MUTED, italic=True)
badge(s, "九年級上學期",
      Inches(0.55), Inches(4.12), w=Inches(1.9), fill=C_AMBER, tc=C_DARK)
multi_para(s, ["富春", "山水", "書信"],
           Inches(9.3), Inches(2.2), Inches(3.7), Inches(2.8),
           size=Pt(30), bold=True,
           color=RGBColor(0x2A, 0x50, 0x45),
           align=PP_ALIGN.CENTER, spacing=Pt(4))
rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_AMBER)

# ─────────────────────────────────────────────────────
# S2  作者介紹
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "作者介紹")

rect(s, Inches(0.55), CT, Inches(3.9), Inches(1.05), fill=C_RIVER)
txbox(s, "吳均",
      Inches(0.55), CT + Inches(0.08), Inches(3.9), Inches(0.55),
      size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Wú Jūn  469–520 CE",
      Inches(0.55), CT + Inches(0.60), Inches(3.9), Inches(0.38),
      size=Pt(12), color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

for i, tag in enumerate(["南朝梁", "文學家", "書信名家"]):
    badge(s, tag, Inches(0.55 + i * 1.38), CT + Inches(1.22),
          w=Inches(1.25), fill=C_MIST)

items = [
    "▸  字叔庠，吳興故鄣（今浙江安吉）人。南朝梁時文學家、史學家，"
    "才學出衆，擅長寫景狀物，文筆清新秀麗，被稱為「吳均體」",
    "▸  本文是寫給友人宋元思的一封書信，描繪作者乘船遊覽富陽至桐廬一帶的山水，"
    "借奇山異水抒發超然脫俗、嚮往自然的情懷",
    "▸  全文僅一百餘字，卻以精煉的語言寫盡富春江的水色山光，"
    "句句工整對仗，被後代文學家視為「小品書信文」的典範作品",
]
multi_para(s, items,
           Inches(0.65), CT + Inches(1.82), Inches(12.0), Inches(4.0),
           size=Pt(16), color=C_GRAY, spacing=Pt(14))

# ─────────────────────────────────────────────────────
# S3  背景知識
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "背景知識：富春江與書信文體")

for i, (icon, head, body) in enumerate([
    ("🏞", "富春江的地理",
     "富春江位於今浙江省境內，是錢塘江的上游。"
     "從富陽至桐廬的江段，兩岸奇峰連綿、江水碧綠清澈，"
     "自古即有「奇山異水，天下獨絕」的美譽，"
     "是歷代文人墨客必訪之地。"),
    ("✉", "書信體古文",
     "書信（書）是古代私人往來的文字，屬應用文體。"
     "吳均以書信形式用清麗駢文描繪山水，寓情於景，"
     "是借「遊記書信」抒發出世之志的典型範例。"
     "文中幾乎每句皆對仗工整，屬「駢文」風格。"),
    ("🖌", "駢文與對偶",
     "「駢文」是講究句式對仗、音律工整的文體。"
     "「水皆縹碧，千丈見底；游魚細石，直視無礙」"
     "是典型對偶句。本文對偶句貫穿全篇，"
     "是學習對偶修辭的絕佳範文。"),
]):
    col = i * Inches(4.3) + Inches(0.50)
    rect(s, col, CT + Inches(0.05), Inches(4.05), Inches(5.5),
         fill=C_WHITE, line_color=C_MIST, line_w=Pt(1.5))
    txbox(s, icon, col + Inches(0.15), CT + Inches(0.12),
          Inches(0.8), Inches(0.7), size=Pt(28), align=PP_ALIGN.CENTER)
    txbox(s, head, col + Inches(0.15), CT + Inches(0.85),
          Inches(3.75), Inches(0.55),
          size=Pt(17), bold=True, color=C_RIVER)
    txbox(s, body, col + Inches(0.15), CT + Inches(1.48),
          Inches(3.75), Inches(3.9),
          size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S4–S8  課文全讀（各段各頁，20pt）
# ─────────────────────────────────────────────────────
slides_data = [
    ("課文全文（一）", "序段",   P1, T1, C_MIST),
    ("課文全文（二）", "水之奇", P2, T2, C_MIST),
    ("課文全文（三）", "山之奇", P3, T3, C_MOUNTAIN),
    ("課文全文（四）", "聲與境", P4, T4, C_RIVER),
    ("課文全文（五）", "尾段",   P5, T5, C_AMBER),
]

for title, label, text, trans, accent in slides_data:
    s = new_slide(prs)
    bg(s, C_CREAM)
    strip_top(s); strip_bot(s)
    slide_title(s, title)
    full_read_bilingual(s, label, text, trans, accent=accent)

# ─────────────────────────────────────────────────────
# S9  對偶手法賞析
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "對偶手法賞析")

pairs = [
    ("水皆縹碧，千丈見底",       "游魚細石，直視無礙"),
    ("急湍甚箭，猛浪若奔",       "負勢競上，互相軒邈"),
    ("泉水激石，泠泠作響",       "好鳥相鳴，嚶嚶成韻"),
    ("蟬則千轉不窮",             "猨則百叫無絕"),
    ("鳶飛戾天者，望峰息心",     "經綸世務者，窺谷忘反"),
]

N_PAIRS  = len(pairs)
GAP_P    = Inches(0.04)
BH_P     = (CH - Inches(0.08) - (N_PAIRS - 1) * GAP_P) // N_PAIRS
CARD_H   = BH_P
TEXT_H   = Inches(0.44)
V_OFF    = int((CARD_H - TEXT_H) / 2)
HALF_W   = Inches(5.72)
ARROW_W  = Inches(0.60)
ARROW_X  = Inches(0.45) + HALF_W + Inches(0.07)
RIGHT_X  = ARROW_X + ARROW_W + Inches(0.07)

for i, (left, right) in enumerate(pairs):
    t = CT + Inches(0.04) + i * (CARD_H + GAP_P)
    rect(s, Inches(0.45), t, HALF_W, CARD_H,
         fill=C_PANEL, line_color=C_MIST, line_w=Pt(1))
    txbox(s, left, Inches(0.60), t + V_OFF, HALF_W - Inches(0.20), TEXT_H,
          size=Pt(17), color=C_RIVER, bold=True)
    txbox(s, "↔", ARROW_X, t + V_OFF, ARROW_W, TEXT_H,
          size=Pt(16), color=C_AMBER, bold=True, align=PP_ALIGN.CENTER)
    rect(s, RIGHT_X, t, HALF_W, CARD_H,
         fill=C_PANEL, line_color=C_MOUNTAIN, line_w=Pt(1))
    txbox(s, right, RIGHT_X + Inches(0.15), t + V_OFF,
          HALF_W - Inches(0.20), TEXT_H,
          size=Pt(17), color=C_MOUNTAIN, bold=True)

# ─────────────────────────────────────────────────────
# S10  文章結構分析
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "文章結構分析")

sections = [
    (C_MIST,     "序段",   "風煙俱淨……天下獨絕",    "總覽：奇山異水，開門見山"),
    (C_MIST,     "水之奇", "水皆縹碧……猛浪若奔",    "靜水清澈 → 急湍如箭，動靜對比"),
    (C_MOUNTAIN, "山之奇", "夾峰高山……千百成峰",    "山木寒挺 → 爭高萬峰，氣勢磅礴"),
    (C_RIVER,    "聲與境", "泉水激石……窺谷忘反",    "泉聲鳥鳴蟬猿 → 出世感悟"),
    (C_AMBER,    "尾段",   "橫柯上蔽……有時見日",    "餘韻：光影交錯，悠然收束"),
]

N_SEC  = len(sections)
GAP_S  = Inches(0.05)
BH_S   = (CH - Inches(0.05) - (N_SEC - 1) * GAP_S) // N_SEC
TEXT_H2 = Inches(0.42)
V_OFF2  = int((BH_S - TEXT_H2) / 2)
LABEL_W = Inches(1.50)
ORIG_X  = Inches(0.45) + LABEL_W + Inches(0.10)
ORIG_W  = Inches(5.40)
CMT_X   = ORIG_X + ORIG_W + Inches(0.10)
CMT_W   = W - CMT_X - Inches(0.40)

for i, (hc, label, orig, comment) in enumerate(sections):
    t = CT + Inches(0.05) + i * (BH_S + GAP_S)
    rect(s, Inches(0.45), t, LABEL_W, BH_S, fill=hc)
    txbox(s, label, Inches(0.45), t + V_OFF2, LABEL_W, TEXT_H2,
          size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, ORIG_X, t, ORIG_W, BH_S,
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    txbox(s, orig, ORIG_X + Inches(0.15), t + V_OFF2,
          ORIG_W - Inches(0.20), TEXT_H2, size=Pt(15), color=C_DARK)
    rect(s, CMT_X, t, CMT_W, BH_S,
         fill=C_PANEL, line_color=C_LGRAY, line_w=Pt(1))
    txbox(s, comment, CMT_X + Inches(0.15), t + V_OFF2,
          CMT_W - Inches(0.20), TEXT_H2, size=Pt(14.5), color=C_GRAY)

# ─────────────────────────────────────────────────────
# S11  重要詞彙（一）：縹碧 / 泠泠 / 嚶嚶
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "重要詞彙（一）")

vocab_a = [
    ("縹碧", "ㄆㄧㄠˇ ㄅㄧˋ",
     "青白透亮的顏色，形容水色清澈而帶有藍綠光澤",
     "水皆縹碧，千丈見底", C_MIST),
    ("泠泠", "ㄌㄧㄥˊ ㄌㄧㄥˊ",
     "泉水或水流清澈悅耳的聲音，形容聲音輕脆清亮",
     "泉水激石，泠泠作響", C_MIST),
    ("嚶嚶", "ㄧㄥ ㄧㄥ",
     "鳥鳴聲輕細婉轉，相互呼應的鳴唱聲",
     "好鳥相鳴，嚶嚶成韻", C_MOUNTAIN),
]

COLS   = 3
COL_W  = (W - Inches(0.90)) / COLS
ROW_H  = CH - Inches(0.08)
GAP_C  = Inches(0.06)

for idx, (word, phonetic, defn, ex, hc) in enumerate(vocab_a):
    l = Inches(0.45) + idx * (COL_W + GAP_C)
    t = CT + Inches(0.04)
    w = COL_W
    h = ROW_H
    rect(s, l, t, w, h, fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, w, Inches(0.80), fill=hc)
    txbox(s, word,
          l + Inches(0.12), t + Inches(0.06),
          Inches(1.8), Inches(0.65),
          size=Pt(26), bold=True, color=C_WHITE)
    txbox(s, phonetic,
          l + Inches(1.95), t + Inches(0.22),
          Inches(2.0), Inches(0.40),
          size=Pt(12), color=RGBColor(0xC0, 0xE0, 0xFF), italic=True)
    rect(s, l + Inches(0.12), t + Inches(0.88),
         w - Inches(0.24), Inches(0.03), fill=C_LGRAY)
    txbox(s, defn,
          l + Inches(0.12), t + Inches(0.96),
          w - Inches(0.24), h - Inches(1.72),
          size=Pt(14), color=C_DARK)
    rect(s, l + Inches(0.12), t + h - Inches(0.68),
         w - Inches(0.24), Inches(0.58), fill=C_PANEL)
    txbox(s, "例：" + ex,
          l + Inches(0.18), t + h - Inches(0.65),
          w - Inches(0.36), Inches(0.54),
          size=Pt(12), color=C_MOUNTAIN, italic=True)

# ─────────────────────────────────────────────────────
# S12  重要詞彙（二）：軒邈 / 戾天 / 經綸
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "重要詞彙（二）")

vocab_b = [
    ("軒邈", "ㄒㄩㄢ ㄇㄧㄠˇ",
     "高聳遠遠，形容山峰互相爭比高遠的氣勢。軒：高；邈：遠",
     "負勢競上，互相軒邈", C_MOUNTAIN),
    ("戾天", "ㄌㄧˋ ㄊㄧㄢ",
     "到達天際。戾：到達、抵達。形容鳶鳥飛翔至高空",
     "鳶飛戾天者，望峰息心", C_RIVER),
    ("經綸", "ㄐㄧㄥ ㄌㄨㄣˊ",
     "籌劃處理政務的事務。經：治理；綸：整理。此處指一切仕途俗務",
     "經綸世務者，窺谷忘反", C_RIVER),
]

for idx, (word, phonetic, defn, ex, hc) in enumerate(vocab_b):
    l = Inches(0.45) + idx * (COL_W + GAP_C)
    t = CT + Inches(0.04)
    w = COL_W
    h = ROW_H
    rect(s, l, t, w, h, fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, w, Inches(0.80), fill=hc)
    txbox(s, word,
          l + Inches(0.12), t + Inches(0.06),
          Inches(1.8), Inches(0.65),
          size=Pt(26), bold=True, color=C_WHITE)
    txbox(s, phonetic,
          l + Inches(1.95), t + Inches(0.22),
          Inches(2.0), Inches(0.40),
          size=Pt(12), color=RGBColor(0xC0, 0xE0, 0xFF), italic=True)
    rect(s, l + Inches(0.12), t + Inches(0.88),
         w - Inches(0.24), Inches(0.03), fill=C_LGRAY)
    txbox(s, defn,
          l + Inches(0.12), t + Inches(0.96),
          w - Inches(0.24), h - Inches(1.72),
          size=Pt(14), color=C_DARK)
    rect(s, l + Inches(0.12), t + h - Inches(0.68),
         w - Inches(0.24), Inches(0.58), fill=C_PANEL)
    txbox(s, "例：" + ex,
          l + Inches(0.18), t + h - Inches(0.65),
          w - Inches(0.36), Inches(0.54),
          size=Pt(12), color=C_MOUNTAIN, italic=True)

# ─────────────────────────────────────────────────────
# S13  主題分析
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "主題分析")

for i, (head, hcolor, items) in enumerate([
    ("自然之美", C_MIST,
     ["靜水如鏡：縹碧千丈，游魚細石直視無礙",
      "急湍如箭：動靜對比，山水各具性格",
      "聲景交融：泉聲、鳥鳴、蟬聲、猿聲",
      "文字精煉簡潔，卻能召喚立體畫面",
      "「奇山異水，天下獨絕」是全文主題"]),
    ("出世情懷", C_MOUNTAIN,
     ["「望峰息心」——名利之心頓消",
      "「窺谷忘反」——俗世羈絆皆忘",
      "吳均藉山水景色，表達不戀仕途",
      "以「境」象徵「心」是古文常用技法",
      "類似陶淵明歸田園的精神取向"]),
    ("動靜對比", C_AMBER,
     ["靜：縹碧千丈，游魚細石——透明靜謐",
      "動：急湍甚箭，猛浪若奔——磅礴洶湧",
      "靜：橫柯上蔽，在晝猶昏——林蔭幽深",
      "動：疎條交映，有時見日——光影流動",
      "動靜交替，讓文字充滿生命節奏"]),
]):
    col = i * Inches(4.25) + Inches(0.45)
    rect(s, col, CT + Inches(0.05), Inches(4.0), Inches(5.5),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, col, CT + Inches(0.05), Inches(4.0), Inches(0.72), fill=hcolor)
    txbox(s, head, col + Inches(0.12), CT + Inches(0.10),
          Inches(3.75), Inches(0.62),
          size=Pt(19), bold=True, color=C_WHITE)
    multi_para(s, ["▸  " + it for it in items],
               col + Inches(0.18), CT + Inches(0.92),
               Inches(3.7), Inches(4.5),
               size=Pt(15), color=C_GRAY, spacing=Pt(10))

# ─────────────────────────────────────────────────────
# S14  核心金句
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)
rect(s, 0, 0, W, H, fill=RGBColor(0x07, 0x10, 0x1A))

txbox(s, "核心金句",
      Inches(0.6), Inches(0.28), Inches(5.0), Inches(0.55),
      size=Pt(16), bold=True, color=C_MIST)
rect(s, Inches(0.6), Inches(0.80), Inches(8.0), Inches(0.05), fill=C_AMBER)

for i, (q, qc, sz) in enumerate([
    ("「奇山異水，天下獨絕。」",
     C_WHITE, Pt(22)),
    ("「急湍甚箭，猛浪若奔。\n"
     "  蟬則千轉不窮，猨則百叫無絕。」",
     RGBColor(0x9A, 0xD4, 0xC4), Pt(19)),
    ("「鳶飛戾天者，望峰息心；\n"
     "  經綸世務者，窺谷忘反。」",
     C_AMBER, Pt(19)),
]):
    t = Inches(1.05) + i * Inches(2.02)
    rect(s, Inches(0.55), t + Inches(0.12),
         Inches(0.06), Inches(1.65),
         fill=C_AMBER if i == 2 else C_MIST)
    txbox(s, q, Inches(0.75), t,
          Inches(12.3), Inches(1.85),
          size=sz, bold=(i == 0), color=qc, italic=True)

# ─────────────────────────────────────────────────────
# S15  課堂討論
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_CREAM)
strip_top(s); strip_bot(s)
slide_title(s, "課堂討論")

for i, (qno, qtxt) in enumerate([
    ("Q1", "課文開頭說「奇山異水，天下獨絕」，\n作者用哪些描寫證明「奇」與「絕」？"),
    ("Q2", "「急湍甚箭，猛浪若奔」與\n「游魚細石，直視無礙」形成什麼樣的對比？"),
    ("Q3", "「鳶飛戾天者，望峰息心；經綸世務者，窺谷忘反」——\n這兩句暗示了作者對什麼樣生活的嚮往？"),
    ("Q4", "本文幾乎每句都對仗工整，\n試找出你最喜歡的一組對偶句，說明欣賞原因。"),
]):
    row, col_i = i // 2, i % 2
    l = Inches(0.45) + col_i * Inches(6.4)
    t = CT + Inches(0.08) + row * Inches(2.9)
    rect(s, l, t, Inches(6.1), Inches(2.72),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, l, t, Inches(1.0), Inches(2.72), fill=C_RIVER)
    txbox(s, qno, l + Inches(0.05), t + Inches(0.95),
          Inches(0.88), Inches(0.82),
          size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txbox(s, qtxt, l + Inches(1.12), t + Inches(0.20),
          Inches(4.82), Inches(2.32),
          size=Pt(15.5), color=C_DARK)

# ─────────────────────────────────────────────────────
# S16  總結
# ─────────────────────────────────────────────────────
s = new_slide(prs)
bg(s, C_COVER)
rect(s, 0, 0, Inches(5.5), H, fill=RGBColor(0x07, 0x10, 0x1A))
rect(s, Inches(5.5), 0, Inches(0.06), H, fill=C_MIST)

txbox(s, "學習重點總結",
      Inches(0.5), Inches(0.6), Inches(4.5), Inches(0.6),
      size=Pt(16), bold=True, color=C_MIST)
rect(s, Inches(0.5), Inches(1.12), Inches(4.5), Inches(0.05), fill=C_AMBER)

multi_para(s,
    ["吳均", "南朝梁・469–520", "字叔庠，吳興故鄣人", "",
     "與宋元思書", "書信體古文", "第五冊第5課"],
    Inches(0.55), Inches(1.28), Inches(4.5), Inches(5.5),
    size=Pt(16), color=RGBColor(0xB0, 0xD0, 0xC8), spacing=Pt(8))

for i, (head, body) in enumerate([
    ("對偶句貫穿全篇",       "幾乎每句皆對仗，是學習駢文與對偶修辭的最佳範本"),
    ("動靜對比寫山水",       "靜水透澈 vs. 急湍如箭，山高林寒 vs. 光影交映"),
    ("借景抒情的出世志",     "「望峰息心、窺谷忘反」隱含不慕名利的人生態度"),
    ("精煉的書信古文",       "全文一百餘字，內容豐富完整，是古文精煉的典範"),
]):
    t = CT - Inches(0.04) + i * Inches(1.52)
    rect(s, Inches(5.85), t, Inches(7.0), Inches(1.40),
         fill=C_WHITE, line_color=C_LGRAY, line_w=Pt(1))
    rect(s, Inches(5.85), t, Inches(0.28), Inches(1.40), fill=C_MOUNTAIN)
    txbox(s, head, Inches(6.22), t + Inches(0.12),
          Inches(6.45), Inches(0.52),
          size=Pt(17), bold=True, color=C_RIVER)
    txbox(s, body, Inches(6.22), t + Inches(0.65),
          Inches(6.45), Inches(0.62),
          size=Pt(14.5), color=C_GRAY)

rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_AMBER)

# ── 儲存 ─────────────────────────────────────────────
OUT = r"C:\Users\yuchi\OneDrive\Desktop\教學日誌\與宋元思書.pptx"
prs.save(OUT)
print(f"已儲存：{OUT}")
